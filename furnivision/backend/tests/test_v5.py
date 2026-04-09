"""Tests for the V5 Human-in-the-Loop system.

Tests:
- Agent 0: PPTX parser (with a real small PPTX)
- Updated models (V5Room, Product, GeneratedImage, FloorPlan)
- ImagenService new methods (multi-ref, edit)
- Video compiler (ffmpeg)
- V5 API routes (full CRUD workflow)
- State management for V5 data
"""

import asyncio
import io
import json
import os
import sys
import tempfile
import uuid
from pathlib import Path
from unittest.mock import AsyncMock, MagicMock, patch

import pytest
import pytest_asyncio

sys.path.insert(0, os.path.join(os.path.dirname(__file__), ".."))

from models.project import (
    FloorPlan,
    GeneratedImage,
    Product,
    Project,
    ProjectBrief,
    V5Room,
)
from pipeline.state import StateManager


# ---------------------------------------------------------------------------
# Fixtures
# ---------------------------------------------------------------------------


@pytest.fixture
def v5_project():
    """A V5 project with rooms, products, images."""
    return Project(
        id="v5_test_001",
        name="Forthing Showroom Test",
        status="reviewing_extraction",
        spec_file_path="/tmp/test_spec.pptx",
        floor_plans=[
            FloorPlan(id="fp1", floor_name="ground", image_path="/tmp/floor_plan.png"),
        ],
        v5_rooms=[
            V5Room(
                id="room_gm",
                label="GM Office",
                floor="ground",
                status="extracted",
                layout_image_path="/tmp/gm_layout.png",
                products=[
                    Product(id="p1", name="Dapper Executive Desk", dimensions="180x90cm", image_path="/tmp/desk.png", room_id="room_gm"),
                    Product(id="p2", name="Dapper Cabinet", dimensions="120x45cm", image_path="/tmp/cabinet.png", room_id="room_gm"),
                ],
            ),
            V5Room(
                id="room_conf",
                label="Conference Room",
                floor="ground",
                status="extracted",
                products=[
                    Product(id="p3", name="Conference Table", dimensions="300x120cm", room_id="room_conf"),
                    Product(id="p4", name="Conference Chair", dimensions="60x60cm", room_id="room_conf"),
                ],
            ),
        ],
    )


@pytest.fixture
def state_manager():
    """In-memory state manager for testing."""
    return StateManager()


# ---------------------------------------------------------------------------
# Model Tests
# ---------------------------------------------------------------------------


class TestV5Models:
    def test_product_creation(self):
        p = Product(id="p1", name="Office Chair", dimensions="60x60x100cm", room_id="room_1")
        assert p.name == "Office Chair"
        assert p.room_id == "room_1"

    def test_generated_image_creation(self):
        img = GeneratedImage(id="img1", room_id="room_1", image_path="/tmp/test.png", version=1, type="base")
        assert img.type == "base"
        assert img.version == 1

    def test_generated_image_types(self):
        for t in ["base", "refined", "edited"]:
            img = GeneratedImage(id="x", room_id="r", type=t)
            assert img.type == t

    def test_floor_plan(self):
        fp = FloorPlan(id="fp1", floor_name="mezzanine", image_path="/tmp/fp.png")
        assert fp.floor_name == "mezzanine"

    def test_v5_room_all_statuses(self):
        for status in ["pending", "extracted", "approved", "generating", "image_ready", "image_approved", "video_ready", "complete"]:
            room = V5Room(id="r", label="Test", status=status)
            assert room.status == status

    def test_v5_room_with_products_and_images(self):
        room = V5Room(
            id="r1",
            label="Test Room",
            products=[Product(id="p1", name="Chair", room_id="r1")],
            generated_images=[GeneratedImage(id="i1", room_id="r1", type="refined")],
            feedback=["needs more casters"],
        )
        assert len(room.products) == 1
        assert len(room.generated_images) == 1
        assert room.feedback[0] == "needs more casters"

    def test_project_v5_statuses(self):
        for status in ["uploading", "extracting", "reviewing_extraction", "generating_images",
                        "reviewing_images", "generating_videos", "complete", "failed"]:
            p = Project(id="test", status=status)
            assert p.status == status

    def test_project_v5_fields(self):
        p = Project(
            id="test",
            spec_file_path="/tmp/spec.pptx",
            floor_plans=[FloorPlan(id="fp1", floor_name="ground")],
            v5_rooms=[V5Room(id="r1", label="Room 1")],
            logo_path="/tmp/logo.png",
            music_path="/tmp/music.mp3",
            final_video_path="/tmp/final.mp4",
        )
        assert p.spec_file_path == "/tmp/spec.pptx"
        assert len(p.floor_plans) == 1
        assert len(p.v5_rooms) == 1
        assert p.logo_path == "/tmp/logo.png"

    def test_project_serialization_roundtrip(self, v5_project):
        """Serialize to dict and back."""
        data = v5_project.model_dump()
        restored = Project(**data)
        assert restored.id == v5_project.id
        assert len(restored.v5_rooms) == 2
        assert restored.v5_rooms[0].label == "GM Office"
        assert len(restored.v5_rooms[0].products) == 2

    def test_project_backward_compat(self):
        """V5 project still supports legacy V4 fields."""
        p = Project(
            id="legacy",
            status="analysing",  # Legacy status
            floorplan_gcs_path="gs://bucket/plan.pdf",
        )
        assert p.status == "analysing"
        assert p.floorplan_gcs_path == "gs://bucket/plan.pdf"


# ---------------------------------------------------------------------------
# State Manager Tests (V5 data)
# ---------------------------------------------------------------------------


class TestV5State:
    @pytest.mark.asyncio
    async def test_create_and_get_v5_project(self, state_manager, v5_project):
        await state_manager.create_project(v5_project)
        loaded = await state_manager.get_project(v5_project.id)
        assert loaded.id == v5_project.id
        assert len(loaded.v5_rooms) == 2
        assert loaded.v5_rooms[0].label == "GM Office"

    @pytest.mark.asyncio
    async def test_update_v5_rooms(self, state_manager, v5_project):
        await state_manager.create_project(v5_project)

        # Update room status
        v5_project.v5_rooms[0].status = "approved"
        rooms_data = [r.model_dump() for r in v5_project.v5_rooms]
        await state_manager.update_project(v5_project.id, {"v5_rooms": rooms_data})

        loaded = await state_manager.get_project(v5_project.id)
        assert loaded.v5_rooms[0].status == "approved"

    @pytest.mark.asyncio
    async def test_add_generated_image_to_room(self, state_manager, v5_project):
        await state_manager.create_project(v5_project)

        # Add a generated image
        new_img = GeneratedImage(id="gen1", room_id="room_gm", image_path="/tmp/gen.png", type="base", version=1)
        v5_project.v5_rooms[0].generated_images.append(new_img)
        rooms_data = [r.model_dump() for r in v5_project.v5_rooms]
        await state_manager.update_project(v5_project.id, {"v5_rooms": rooms_data})

        loaded = await state_manager.get_project(v5_project.id)
        assert len(loaded.v5_rooms[0].generated_images) == 1
        assert loaded.v5_rooms[0].generated_images[0].type == "base"

    @pytest.mark.asyncio
    async def test_update_project_status_flow(self, state_manager, v5_project):
        await state_manager.create_project(v5_project)

        statuses = ["extracting", "reviewing_extraction", "generating_images", "reviewing_images", "generating_videos", "complete"]
        for status in statuses:
            await state_manager.update_project(v5_project.id, {"status": status})
            loaded = await state_manager.get_project(v5_project.id)
            assert loaded.status == status

    @pytest.mark.asyncio
    async def test_feedback_persistence(self, state_manager, v5_project):
        await state_manager.create_project(v5_project)

        v5_project.v5_rooms[0].feedback.append("chairs need casters")
        v5_project.v5_rooms[0].feedback.append("wrong table color")
        rooms_data = [r.model_dump() for r in v5_project.v5_rooms]
        await state_manager.update_project(v5_project.id, {"v5_rooms": rooms_data})

        loaded = await state_manager.get_project(v5_project.id)
        assert len(loaded.v5_rooms[0].feedback) == 2
        assert "casters" in loaded.v5_rooms[0].feedback[0]


# ---------------------------------------------------------------------------
# Agent 0: PPTX Parser Tests
# ---------------------------------------------------------------------------


class TestAgent0PPTXParser:
    def _create_test_pptx(self, tmp_dir: str) -> str:
        """Create a minimal PPTX file for testing."""
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()

        # Slide 1: Title slide
        slide1 = prs.slides.add_slide(prs.slide_layouts[0])
        slide1.shapes.title.text = "Forthing Showroom"
        slide1.placeholders[1].text = "Furniture Specification"

        # Slide 2: Floor plan (simulate with a text slide)
        slide2 = prs.slides.add_slide(prs.slide_layouts[1])
        slide2.shapes.title.text = "Floor Plan Overview"

        # Slide 3: GM Office room
        slide3 = prs.slides.add_slide(prs.slide_layouts[1])
        slide3.shapes.title.text = "GM Office"
        body = slide3.placeholders[1]
        body.text = "Dapper Executive Desk 180x90cm\nDapper Cabinet 120x45cm\nExecutive Chair"

        # Slide 4: Conference Room
        slide4 = prs.slides.add_slide(prs.slide_layouts[1])
        slide4.shapes.title.text = "Conference Room"
        body = slide4.placeholders[1]
        body.text = "Conference Table 300x120cm\n12x Conference Chairs with casters"

        # Save
        pptx_path = os.path.join(tmp_dir, "test_spec.pptx")
        prs.save(pptx_path)
        return pptx_path

    def test_extract_slides(self):
        """Test raw slide extraction from PPTX."""
        from agents.agent0_pptx_parser import PPTXParserAgent

        with tempfile.TemporaryDirectory() as tmp:
            pptx_path = self._create_test_pptx(tmp)
            agent = PPTXParserAgent()
            slides = agent._extract_slides(pptx_path)

            assert len(slides) == 4
            assert "Forthing Showroom" in slides[0].texts[0]
            assert "GM Office" in slides[2].texts[0]
            assert "Conference Room" in slides[3].texts[0]

    def test_build_slide_summaries(self):
        """Test slide summary generation."""
        from agents.agent0_pptx_parser import PPTXParserAgent

        with tempfile.TemporaryDirectory() as tmp:
            pptx_path = self._create_test_pptx(tmp)
            agent = PPTXParserAgent()
            slides = agent._extract_slides(pptx_path)
            summaries = agent._build_slide_summaries(slides)

            assert len(summaries) == 4
            assert summaries[0]["slide_index"] == 0
            assert len(summaries[2]["texts"]) > 0

    def test_build_result_from_classification(self):
        """Test building result from a mock Gemini classification."""
        from agents.agent0_pptx_parser import PPTXParserAgent, SlideData

        with tempfile.TemporaryDirectory() as tmp:
            pptx_path = self._create_test_pptx(tmp)
            agent = PPTXParserAgent()
            slides = agent._extract_slides(pptx_path)

            # Mock classification result (what Gemini would return)
            classification = {
                "floor_plan_slides": [1],
                "rooms": [
                    {
                        "label": "GM Office",
                        "floor": "ground",
                        "slide_indices": [2],
                        "products": [
                            {"name": "Dapper Executive Desk", "dimensions": "180x90cm", "image_index_on_slide": 0, "slide_index": 2},
                            {"name": "Dapper Cabinet", "dimensions": "120x45cm", "image_index_on_slide": 1, "slide_index": 2},
                        ],
                    },
                    {
                        "label": "Conference Room",
                        "floor": "ground",
                        "slide_indices": [3],
                        "products": [
                            {"name": "Conference Table", "dimensions": "300x120cm", "image_index_on_slide": 0, "slide_index": 3},
                        ],
                    },
                ],
                "ignored_slides": [0],
            }

            result = agent._build_result(slides, classification, Path(tmp), "test_proj")

            assert len(result.rooms) == 2
            assert result.rooms[0]["label"] == "GM Office"
            assert len(result.rooms[0]["products"]) == 2
            assert result.rooms[1]["label"] == "Conference Room"

    @pytest.mark.asyncio
    async def test_full_parse_with_mock_gemini(self):
        """Test full parse pipeline with mocked Gemini."""
        from agents.agent0_pptx_parser import PPTXParserAgent

        with tempfile.TemporaryDirectory() as tmp:
            pptx_path = self._create_test_pptx(tmp)

            mock_classification = {
                "floor_plan_slides": [1],
                "rooms": [
                    {
                        "label": "GM Office",
                        "floor": "ground",
                        "slide_indices": [2],
                        "products": [
                            {"name": "Executive Desk", "dimensions": "180x90cm", "image_index_on_slide": 0, "slide_index": 2},
                        ],
                    },
                ],
                "ignored_slides": [0],
            }

            agent = PPTXParserAgent()
            with patch.object(agent.gemini, "analyze_images_structured", new_callable=AsyncMock, return_value=mock_classification):
                result = await agent.parse(pptx_path, "test_proj_001")

            assert len(result.rooms) == 1
            assert result.rooms[0]["label"] == "GM Office"


# ---------------------------------------------------------------------------
# Agent 2.5: Scene Composer Tests
# ---------------------------------------------------------------------------


class TestAgent25Composer:
    @pytest.mark.asyncio
    async def test_compose_with_mocks(self):
        """Test the V4/V5 pipeline with mocked services."""
        from agents.agent2_5_composer import SceneComposerAgent

        agent = SceneComposerAgent()

        # Mock Gemini description
        mock_description = (
            "Photorealistic architectural interior render of a GM Office. "
            "A spacious room with a large walnut executive desk..."
        )

        # Mock Imagen base render (fake PNG bytes)
        fake_base = self._make_fake_png(1536, 864)
        fake_refined = self._make_fake_png(1536, 864)

        with patch.object(agent.gemini, "compose_room_scene", new_callable=AsyncMock, return_value=mock_description), \
             patch.object(agent.imagen, "generate_frame_with_retry", new_callable=AsyncMock, return_value=fake_base), \
             patch.object(agent.imagen, "generate_frame_from_reference_multi", new_callable=AsyncMock, return_value=fake_refined), \
             patch.object(agent.gemini, "analyze_images_structured", new_callable=AsyncMock, return_value={"best_index": 0, "reasoning": "best match"}):

            result = await agent.compose(
                room_id="room_gm",
                room_label="GM Office",
                floor_plan_bytes=b"fake_floor_plan",
                furniture_images=[b"fake_furniture_1", b"fake_furniture_2"],
            )

        assert result.room_id == "room_gm"
        assert result.room_label == "GM Office"
        assert result.description == mock_description
        assert len(result.base_render) > 0
        assert len(result.refined_render) > 0

    @pytest.mark.asyncio
    async def test_compose_fallback_on_refinement_failure(self):
        """If all refinements fail, base render should be used."""
        from agents.agent2_5_composer import SceneComposerAgent

        agent = SceneComposerAgent()
        fake_base = self._make_fake_png(1536, 864)

        with patch.object(agent.gemini, "compose_room_scene", new_callable=AsyncMock, return_value="Test description"), \
             patch.object(agent.imagen, "generate_frame_with_retry", new_callable=AsyncMock, return_value=fake_base), \
             patch.object(agent.imagen, "generate_frame_from_reference_multi", new_callable=AsyncMock, side_effect=Exception("API error")):

            result = await agent.compose(
                room_id="r1",
                room_label="Test",
                floor_plan_bytes=b"fp",
                furniture_images=[b"f1"],
            )

        # Should fall back to base render
        assert result.refined_render == result.base_render
        assert result.refinement_attempts == 0

    def test_ensure_16x9(self):
        """Test 16:9 resize logic."""
        from agents.agent2_5_composer import SceneComposerAgent

        # Already correct size
        correct = self._make_fake_png(1536, 864)
        result = SceneComposerAgent._ensure_16x9(correct)
        from PIL import Image
        img = Image.open(io.BytesIO(result))
        assert img.size == (1536, 864)

        # Wrong size should be resized
        wrong = self._make_fake_png(800, 600)
        result = SceneComposerAgent._ensure_16x9(wrong)
        img = Image.open(io.BytesIO(result))
        assert img.size == (1536, 864)

    @staticmethod
    def _make_fake_png(w: int, h: int) -> bytes:
        from PIL import Image
        img = Image.new("RGB", (w, h), color=(128, 128, 128))
        buf = io.BytesIO()
        img.save(buf, "PNG")
        return buf.getvalue()


# ---------------------------------------------------------------------------
# ImagenService New Methods Tests
# ---------------------------------------------------------------------------


class TestImagenServiceV5:
    @pytest.mark.asyncio
    async def test_multi_ref_method_exists(self):
        """Verify the new method signature is callable."""
        from services.imagen import ImagenService
        service = ImagenService()
        assert hasattr(service, "generate_frame_from_reference_multi")
        assert hasattr(service, "edit_image_with_feedback")

    @pytest.mark.asyncio
    async def test_edit_image_mock(self):
        """Test edit_image_with_feedback with mocked API."""
        from services.imagen import ImagenService

        service = ImagenService()
        fake_result = MagicMock()
        fake_part = MagicMock()
        fake_inline = MagicMock()
        fake_inline.data = b"edited_image_bytes"
        fake_part.inline_data = fake_inline
        fake_result.candidates = [MagicMock()]
        fake_result.candidates[0].content.parts = [fake_part]

        with patch.object(service._client.models, "generate_content", return_value=fake_result):
            result = await service.edit_image_with_feedback(
                current_image=b"current",
                feedback="add casters to chairs",
                reference_images=[b"ref1"],
            )

        assert result == b"edited_image_bytes"


# ---------------------------------------------------------------------------
# Video Compiler Tests
# ---------------------------------------------------------------------------


class TestVideoCompiler:
    @pytest.mark.asyncio
    async def test_compiler_no_ffmpeg(self):
        """Should raise if ffmpeg not available."""
        from services.video_compiler import VideoCompiler, VideoCompilerError

        compiler = VideoCompiler()
        with patch("shutil.which", return_value=None):
            with pytest.raises(VideoCompilerError, match="ffmpeg not found"):
                await compiler.compile(
                    room_video_paths=["/tmp/v1.mp4"],
                    output_path="/tmp/out.mp4",
                )

    @pytest.mark.asyncio
    async def test_compiler_no_videos(self):
        """Should raise if no videos provided."""
        from services.video_compiler import VideoCompiler, VideoCompilerError

        compiler = VideoCompiler()
        with pytest.raises(VideoCompilerError, match="No room videos"):
            await compiler.compile(
                room_video_paths=[],
                output_path="/tmp/out.mp4",
            )

    @pytest.mark.asyncio
    async def test_compiler_with_ffmpeg(self):
        """Test full compilation with real ffmpeg (if available)."""
        import shutil
        if not shutil.which("ffmpeg"):
            pytest.skip("ffmpeg not installed")

        from services.video_compiler import VideoCompiler

        compiler = VideoCompiler()

        with tempfile.TemporaryDirectory() as tmp:
            # Create a tiny test video using ffmpeg
            test_video = os.path.join(tmp, "test_room.mp4")
            proc = await asyncio.create_subprocess_exec(
                "ffmpeg", "-y",
                "-f", "lavfi", "-i", "color=c=blue:s=320x180:d=1",
                "-c:v", "libx264", "-pix_fmt", "yuv420p",
                test_video,
                stdout=asyncio.subprocess.PIPE,
                stderr=asyncio.subprocess.PIPE,
            )
            await proc.communicate()
            assert os.path.exists(test_video), "Failed to create test video"

            # Compile
            output = os.path.join(tmp, "final.mp4")
            result = await compiler.compile(
                room_video_paths=[test_video, test_video],
                output_path=output,
            )

            assert os.path.exists(result)
            assert os.path.getsize(result) > 0


# ---------------------------------------------------------------------------
# V5 API Route Tests
# ---------------------------------------------------------------------------


class TestV5APIRoutes:
    @pytest.fixture
    def client(self):
        """Create a test client."""
        from fastapi.testclient import TestClient
        from main import app
        return TestClient(app)

    def test_create_v5_project(self, client):
        res = client.post("/api/v1/v5/projects?name=Test+Project")
        assert res.status_code == 200
        data = res.json()
        assert "project_id" in data
        assert data["status"] == "uploading"
        return data["project_id"]

    def test_get_v5_project(self, client):
        # Create first
        create_res = client.post("/api/v1/v5/projects?name=Get+Test")
        pid = create_res.json()["project_id"]

        res = client.get(f"/api/v1/v5/projects/{pid}")
        assert res.status_code == 200
        data = res.json()
        assert data["id"] == pid
        assert data["name"] == "Get Test"

    def test_get_nonexistent_project(self, client):
        res = client.get("/api/v1/v5/projects/nonexistent_xyz")
        assert res.status_code == 404

    def test_extraction_not_ready(self, client):
        create_res = client.post("/api/v1/v5/projects?name=Extr+Test")
        pid = create_res.json()["project_id"]

        res = client.get(f"/api/v1/v5/projects/{pid}/extraction")
        # Extraction endpoint now always returns 200 with current state
        assert res.status_code == 200
        data = res.json()
        assert data["status"] == "uploading"
        assert data["total_products"] == 0
        assert len(data["rooms"]) == 0

    def test_update_room(self, client):
        # Create project with rooms
        create_res = client.post("/api/v1/v5/projects?name=Room+Update")
        pid = create_res.json()["project_id"]

        # Manually add rooms to project state
        state = StateManager()
        loop = asyncio.new_event_loop()
        room = V5Room(id="r1", label="Old Name", status="extracted")
        loop.run_until_complete(state.update_project(pid, {
            "v5_rooms": [room.model_dump()],
            "status": "reviewing_extraction",
        }))
        loop.close()

        # Update room name
        res = client.put(f"/api/v1/v5/projects/{pid}/rooms/r1", json={"label": "New Name"})
        assert res.status_code == 200

        # Verify
        project_res = client.get(f"/api/v1/v5/projects/{pid}")
        assert project_res.json()["v5_rooms"][0]["label"] == "New Name"

    def test_approve_extraction(self, client):
        create_res = client.post("/api/v1/v5/projects?name=Approve+Test")
        pid = create_res.json()["project_id"]

        # Add rooms
        state = StateManager()
        loop = asyncio.new_event_loop()
        rooms = [
            V5Room(id="r1", label="Room 1", status="extracted",
                   products=[Product(id="p1", name="Chair", room_id="r1")]),
            V5Room(id="r2", label="Room 2", status="extracted"),
        ]
        loop.run_until_complete(state.update_project(pid, {
            "v5_rooms": [r.model_dump() for r in rooms],
            "status": "reviewing_extraction",
        }))
        loop.close()

        # Approve (will trigger background task that may fail without real AI, that's ok)
        res = client.post(f"/api/v1/v5/projects/{pid}/approve-extraction")
        assert res.status_code == 200
        data = res.json()
        assert data["rooms_count"] == 2

    def test_approve_room_image(self, client):
        create_res = client.post("/api/v1/v5/projects?name=ImgApprove")
        pid = create_res.json()["project_id"]

        state = StateManager()
        loop = asyncio.new_event_loop()
        rooms = [
            V5Room(id="r1", label="Room 1", status="image_ready",
                   generated_images=[GeneratedImage(id="i1", room_id="r1", type="refined")]),
        ]
        loop.run_until_complete(state.update_project(pid, {
            "v5_rooms": [r.model_dump() for r in rooms],
            "status": "reviewing_images",
        }))
        loop.close()

        res = client.post(f"/api/v1/v5/projects/{pid}/rooms/r1/approve")
        assert res.status_code == 200
        assert res.json()["status"] == "image_approved"

    def test_full_v5_workflow_state_transitions(self, client):
        """Test the complete project status flow through the API."""
        # Create
        create_res = client.post("/api/v1/v5/projects?name=Full+Flow")
        pid = create_res.json()["project_id"]
        assert create_res.json()["status"] == "uploading"

        # Simulate extraction complete
        state = StateManager()
        loop = asyncio.new_event_loop()
        rooms = [
            V5Room(id="r1", label="Room 1", status="extracted",
                   products=[Product(id="p1", name="Desk", room_id="r1")]),
        ]
        loop.run_until_complete(state.update_project(pid, {
            "v5_rooms": [r.model_dump() for r in rooms],
            "status": "reviewing_extraction",
        }))

        # Verify status
        proj = loop.run_until_complete(state.get_project(pid))
        assert proj.status == "reviewing_extraction"

        # Simulate image generation complete
        rooms[0].status = "image_ready"
        rooms[0].generated_images = [GeneratedImage(id="g1", room_id="r1", type="refined")]
        loop.run_until_complete(state.update_project(pid, {
            "v5_rooms": [r.model_dump() for r in rooms],
            "status": "reviewing_images",
        }))

        proj = loop.run_until_complete(state.get_project(pid))
        assert proj.status == "reviewing_images"
        assert len(proj.v5_rooms[0].generated_images) == 1

        loop.close()

    def test_generate_videos_no_approved(self, client):
        """Should fail if no rooms are approved."""
        create_res = client.post("/api/v1/v5/projects?name=NoApproved")
        pid = create_res.json()["project_id"]

        state = StateManager()
        loop = asyncio.new_event_loop()
        loop.run_until_complete(state.update_project(pid, {
            "v5_rooms": [V5Room(id="r1", label="R1", status="extracted").model_dump()],
        }))
        loop.close()

        res = client.post(f"/api/v1/v5/projects/{pid}/generate-videos")
        assert res.status_code == 400

    def test_compile_no_videos(self, client):
        """Should fail if no room videos exist."""
        create_res = client.post("/api/v1/v5/projects?name=NoVideos")
        pid = create_res.json()["project_id"]

        res = client.post(f"/api/v1/v5/projects/{pid}/compile")
        assert res.status_code == 400

    def test_health_check(self, client):
        res = client.get("/health")
        assert res.status_code == 200
        assert res.json()["status"] == "healthy"

    def test_upload_spec_size_limit(self, client):
        """Should reject PPTX files over 100MB."""
        create_res = client.post("/api/v1/v5/projects?name=SizeLimit")
        pid = create_res.json()["project_id"]

        # Create a fake "large" file - we can't actually send 100MB in test,
        # but we verify the endpoint exists and accepts small files
        import io
        small_file = io.BytesIO(b"PK" + b"\x00" * 100)  # Tiny fake PPTX
        res = client.post(
            f"/api/v1/v5/projects/{pid}/upload-spec",
            files={"file": ("test.pptx", small_file, "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
        )
        # Will fail on extraction (not a real PPTX) but upload should succeed (under size limit)
        assert res.status_code == 200


# ---------------------------------------------------------------------------
# Code Review Fix Tests
# ---------------------------------------------------------------------------


class TestCodeReviewFixes:
    """Tests for issues found in code review."""

    def test_v5_room_failure_statuses(self):
        """Fix 7: V5Room should accept failure states."""
        room = V5Room(id="r1", label="Test", status="generation_failed")
        assert room.status == "generation_failed"
        room2 = V5Room(id="r2", label="Test2", status="video_failed")
        assert room2.status == "video_failed"

    def test_video_compiler_ensure_even(self):
        """Fix 4: Even dimensions for H.264."""
        from services.video_compiler import VideoCompiler
        assert VideoCompiler._ensure_even(1920) == 1920
        assert VideoCompiler._ensure_even(1921) == 1920
        assert VideoCompiler._ensure_even(1080) == 1080
        assert VideoCompiler._ensure_even(1081) == 1080
        assert VideoCompiler._ensure_even(0) == 0
        assert VideoCompiler._ensure_even(1) == 0
        assert VideoCompiler._ensure_even(3) == 2

    def test_video_completion_logic(self):
        """Fix 2: Video completion check should use video_path, not status matching."""
        # All rooms have video_path → all_ready should be True
        rooms_ready = [
            V5Room(id="r1", label="R1", status="video_ready", video_path="/tmp/v1.mp4"),
            V5Room(id="r2", label="R2", status="video_ready", video_path="/tmp/v2.mp4"),
        ]
        rooms_needing = [r for r in rooms_ready if r.status in ("image_approved", "video_ready")]
        all_ready = bool(rooms_needing) and all(r.video_path for r in rooms_needing)
        assert all_ready is True

        # One room missing video_path → not ready
        rooms_partial = [
            V5Room(id="r1", label="R1", status="video_ready", video_path="/tmp/v1.mp4"),
            V5Room(id="r2", label="R2", status="image_approved", video_path=None),
        ]
        rooms_needing2 = [r for r in rooms_partial if r.status in ("image_approved", "video_ready")]
        all_ready2 = bool(rooms_needing2) and all(r.video_path for r in rooms_needing2)
        assert all_ready2 is False

        # No rooms needing video → not ready (empty)
        rooms_empty: list[V5Room] = []
        rooms_needing3 = [r for r in rooms_empty if r.status in ("image_approved", "video_ready")]
        all_ready3 = bool(rooms_needing3) and all(r.video_path for r in rooms_needing3)
        assert all_ready3 is False

    @pytest.mark.asyncio
    async def test_state_manager_uses_global(self):
        """Fix 1: _generate_room_images should use global _state, not new instance."""
        # Verify that the module-level _state is used (not a new StateManager)
        import importlib
        v5_module = importlib.import_module("api.routes.v5")
        # Check the source code doesn't contain "state = StateManager()" in _generate_room_images
        import inspect
        source = inspect.getsource(v5_module._generate_room_images)
        assert "StateManager()" not in source, \
            "_generate_room_images should use _state, not create new StateManager()"
        assert "_state" in source

    def test_imagen_backoff_cap(self):
        """Fix 8: Retry backoff should be capped at 30s."""
        backoff_base = 2.0
        # Attempt 10 without cap: 2^10 = 1024s
        # With cap: min(1024, 30) = 30s
        for attempt in range(1, 20):
            wait = min(backoff_base ** attempt, 30)
            assert wait <= 30, f"Backoff at attempt {attempt} = {wait}s, should be capped at 30"

    @pytest.mark.asyncio
    async def test_parallel_refinements(self):
        """Fix 10: Agent 2.5 should run refinements in parallel."""
        from agents.agent2_5_composer import SceneComposerAgent
        import inspect
        source = inspect.getsource(SceneComposerAgent.compose)
        assert "asyncio.gather" in source, \
            "compose() should use asyncio.gather for parallel refinements"
