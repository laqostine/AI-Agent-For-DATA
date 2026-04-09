"""FurniVision AI — Agent 0: PPTX Spec Parser.

Extracts structured room-by-room data from a furniture specification PPTX file
(like FORTHING SHOWROOM.pptx).

Given:
  - A PPTX file path (spec document with room slides)

Produces:
  - Per-room: label, layout images, product images with names/dimensions
  - Floor plan images from overview slides
  - Product-to-room mapping based on slide structure

Uses python-pptx for image extraction and Gemini for understanding slide
structure (which images are products vs layouts, room names, etc.).
"""

import asyncio
import io
import json
import logging
import re
import uuid
from pathlib import Path

from pptx import Presentation
from PIL import Image

from services.gemini import GeminiService
from config import TEMP_DIR

logger = logging.getLogger(__name__)

_JSON_FENCE_RE = re.compile(r"```(?:json)?\s*\n?(.*?)```", re.DOTALL)


class SlideData:
    """Raw data extracted from a single PPTX slide."""

    def __init__(self, index: int) -> None:
        self.index = index
        self.texts: list[str] = []
        self.images: list[bytes] = []
        self.image_content_types: list[str] = []


class PPTXExtractionResult:
    """Result of PPTX parsing: rooms, products, floor plans."""

    def __init__(self) -> None:
        self.rooms: list[dict] = []
        self.floor_plans: list[dict] = []
        self.all_product_images: list[dict] = []


class PPTXParserAgent:
    """Agent 0 — Parse a PPTX spec document into rooms + products."""

    def __init__(self) -> None:
        self.gemini = GeminiService()

    async def parse(self, pptx_path: str, project_id: str) -> PPTXExtractionResult:
        """Parse a PPTX file and extract rooms, products, and floor plans."""
        logger.info("PPTXParserAgent.parse — file=%s, project=%s", pptx_path, project_id)

        output_dir = TEMP_DIR / f"pptx_extract_{project_id}"
        output_dir.mkdir(parents=True, exist_ok=True)

        # Step 1: Extract raw slide data (text + images via blipFill)
        slides = self._extract_slides(pptx_path)
        logger.info("Extracted %d slides, total images: %d",
                     len(slides), sum(len(s.images) for s in slides))

        # Step 2: Build summaries for Gemini
        slide_summaries = self._build_slide_summaries(slides)

        # Step 3: Gemini classifies slides and maps products to rooms
        classification = await self._classify_slides_with_gemini(slides, slide_summaries)
        logger.info("Gemini classification: %d rooms",
                     len(classification.get("rooms", [])))

        # Step 4: Save images and build result
        result = self._build_result(slides, classification, output_dir, project_id)

        logger.info(
            "PPTX parsing complete: %d rooms, %d products, %d floor plans",
            len(result.rooms), len(result.all_product_images), len(result.floor_plans),
        )
        return result

    # ------------------------------------------------------------------
    # Step 1: Extract raw data from PPTX
    # ------------------------------------------------------------------

    def _extract_slides(self, pptx_path: str) -> list[SlideData]:
        """Extract text and images from every slide.

        Images are found via blipFill references in the XML (handles Freeform,
        AutoShape, Group, and Picture shapes — not just MSO_SHAPE_TYPE.PICTURE).
        """
        prs = Presentation(pptx_path)
        slides: list[SlideData] = []

        for idx, slide in enumerate(prs.slides):
            data = SlideData(index=idx)

            # Extract text from all shapes (including groups)
            self._extract_text_recursive(slide.shapes, data)

            # Extract images via relationship IDs in the XML
            self._extract_images_from_slide(slide, data)

            slides.append(data)

        return slides

    def _extract_text_recursive(self, shapes, data: SlideData) -> None:
        """Recursively extract text from shapes and group shapes."""
        for shape in shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        data.texts.append(text)
            # Recurse into group shapes
            if hasattr(shape, "shapes"):
                self._extract_text_recursive(shape.shapes, data)

    def _extract_images_from_slide(self, slide, data: SlideData) -> None:
        """Extract all embedded images from a slide via blip references.

        Scans the slide XML for r:embed references pointing to image parts.
        This handles all shape types (Freeform, AutoShape, Group, Picture).
        """
        seen_rids: set[str] = set()

        for shape in slide.shapes:
            self._extract_blip_images(shape, slide, data, seen_rids)
            # Also check group children
            if hasattr(shape, "shapes"):
                for sub_shape in shape.shapes:
                    self._extract_blip_images(sub_shape, slide, data, seen_rids)

    def _extract_blip_images(self, shape, slide, data: SlideData, seen_rids: set) -> None:
        """Extract image blobs from a shape's XML blip references."""
        xml = shape._element.xml
        rids = re.findall(r'r:embed="(rId\d+)"', xml)

        for rid in rids:
            if rid in seen_rids:
                continue
            try:
                rel = slide.part.rels[rid]
                part = rel.target_part
                ct = part.content_type or ""
                if ct.startswith("image/"):
                    blob = part.blob
                    if blob and len(blob) > 1000:  # Skip tiny icons/logos
                        data.images.append(blob)
                        data.image_content_types.append(ct)
                        seen_rids.add(rid)
            except Exception as exc:
                logger.debug("Could not extract image rId %s: %s", rid, exc)

    # ------------------------------------------------------------------
    # Step 2: Summaries
    # ------------------------------------------------------------------

    def _build_slide_summaries(self, slides: list[SlideData]) -> list[dict]:
        """Build a text summary of each slide for Gemini."""
        return [
            {
                "slide_index": s.index,
                "texts": s.texts[:20],
                "image_count": len(s.images),
            }
            for s in slides
        ]

    # ------------------------------------------------------------------
    # Step 3: Gemini classification
    # ------------------------------------------------------------------

    async def _classify_slides_with_gemini(
        self,
        slides: list[SlideData],
        summaries: list[dict],
    ) -> dict:
        """Send slide summaries + key images to Gemini for classification.

        Sends the first image from selected slides (room headers + a few products)
        to help Gemini understand the structure.
        """
        # Select representative images: one from every ~3 slides, max 15
        image_parts: list[bytes] = []
        image_slide_map: list[int] = []

        for s in slides:
            if s.images and len(image_parts) < 15:
                # Pick slides that look like room headers or products
                image_parts.append(s.images[0])
                image_slide_map.append(s.index)

        system_prompt = (
            "You are an expert at analyzing furniture specification PPTX documents. "
            "These documents describe room-by-room furniture layouts for commercial "
            "spaces like showrooms, offices, and dealerships.\n\n"
            "The PPTX has this pattern:\n"
            "- Room header slides: contain the room name (e.g. 'ACCUEIL', 'SALES LOUNGE', 'GM OFFICE')\n"
            "  followed by product slides showing individual furniture pieces.\n"
            "- Product slides: contain the product name (e.g. 'TASK CHAIR', 'DESK'), dimensions, and a product image.\n"
            "- Floor plan slides: show overview layouts (e.g. 'GROUND FLOOR', 'Mezzanine Floor').\n"
            "- Some slides may be title/intro slides to ignore.\n\n"
            "You will receive slide text summaries and some representative images."
        )

        user_prompt = (
            f"SLIDE SUMMARIES (64 slides):\n{json.dumps(summaries, indent=2)}\n\n"
            f"IMAGE-TO-SLIDE MAP: Representative images are from slides {image_slide_map}\n\n"
            "Analyze the slide structure and respond with JSON:\n"
            "{\n"
            '  "floor_plan_slides": [<slide indices showing floor plans/overviews>],\n'
            '  "rooms": [\n'
            "    {\n"
            '      "label": "<room name>",\n'
            '      "floor": "<ground or mezzanine>",\n'
            '      "header_slide": <slide index of room header>,\n'
            '      "product_slides": [<slide indices of this room\'s products>],\n'
            '      "products": [\n'
            "        {\n"
            '          "name": "<product name from slide text>",\n'
            '          "dimensions": "<dimensions if mentioned, e.g. 200x90x110>",\n'
            '          "slide_index": <which slide>\n'
            "        }\n"
            "      ]\n"
            "    }\n"
            "  ],\n"
            '  "ignored_slides": [<title/intro/end slides to skip>]\n'
            "}\n\n"
            "IMPORTANT:\n"
            "- Each product slide has exactly ONE product and ONE image\n"
            "- Room header slides have the room name in ALL CAPS\n"
            "- Products after a room header belong to that room until the next header\n"
            "- Floor plan slides show 'GROUND FLOOR' or 'Mezzanine Floor'\n"
            "- Include ALL products you can identify\n"
            "Respond ONLY with JSON. No extra text."
        )

        result = await self.gemini.analyze_images_structured(
            images=image_parts,
            system_prompt=system_prompt,
            user_prompt=user_prompt,
            max_retries=3,
        )
        return result

    # ------------------------------------------------------------------
    # Step 4: Build result
    # ------------------------------------------------------------------

    def _build_result(
        self,
        slides: list[SlideData],
        classification: dict,
        output_dir: Path,
        project_id: str,
    ) -> PPTXExtractionResult:
        """Save extracted images and build the final result."""
        result = PPTXExtractionResult()

        # Floor plans
        for slide_idx in classification.get("floor_plan_slides", []):
            if slide_idx < len(slides) and slides[slide_idx].images:
                fp_id = str(uuid.uuid4())[:8]
                fp_path = output_dir / f"floor_plan_{slide_idx}.png"
                self._save_image(slides[slide_idx].images[0], fp_path)
                result.floor_plans.append({
                    "id": fp_id,
                    "floor_name": "mezzanine" if any(
                        "mezzanine" in t.lower() for t in slides[slide_idx].texts
                    ) else "ground",
                    "image_path": str(fp_path),
                })

        # Rooms and products
        for room_data in classification.get("rooms", []):
            room_id = str(uuid.uuid4())[:8]
            room_label = room_data.get("label", f"Room {room_id}")
            room_floor = room_data.get("floor", "ground")

            # Save room header image as layout
            header_slide = room_data.get("header_slide")
            layout_path = ""
            if header_slide is not None and header_slide < len(slides):
                slide = slides[header_slide]
                if slide.images:
                    layout_path = str(output_dir / f"room_{room_id}_layout.png")
                    self._save_image(slide.images[0], Path(layout_path))

            # Extract products
            products = []
            for prod_data in room_data.get("products", []):
                prod_id = str(uuid.uuid4())[:8]
                prod_name = prod_data.get("name", "Unknown Product")
                prod_dims = prod_data.get("dimensions", "")
                prod_slide = prod_data.get("slide_index", -1)

                # Save product image (each product slide has 1 image)
                prod_image_path = ""
                if 0 <= prod_slide < len(slides) and slides[prod_slide].images:
                    safe_name = re.sub(r'[^\w\-]', '_', prod_name[:25])
                    prod_image_path = str(
                        output_dir / f"product_{prod_id}_{safe_name}.png"
                    )
                    self._save_image(slides[prod_slide].images[0], Path(prod_image_path))

                product = {
                    "id": prod_id,
                    "name": prod_name,
                    "dimensions": prod_dims,
                    "image_path": prod_image_path,
                    "slide_index": prod_slide,
                    "room_id": room_id,
                }
                products.append(product)
                result.all_product_images.append(product)

            result.rooms.append({
                "id": room_id,
                "label": room_label,
                "floor": room_floor,
                "slide_index": header_slide if header_slide is not None else 0,
                "layout_image_path": layout_path,
                "products": products,
            })

        return result

    # ------------------------------------------------------------------
    # Image helpers
    # ------------------------------------------------------------------

    @staticmethod
    def _save_image(image_bytes: bytes, output_path: Path) -> None:
        """Save image bytes to disk, converting to PNG if needed."""
        try:
            img = Image.open(io.BytesIO(image_bytes))
            # Convert CMYK or palette to RGB
            if img.mode in ("CMYK", "P", "LA", "PA"):
                img = img.convert("RGBA" if img.mode in ("LA", "PA") else "RGB")
            img.save(str(output_path), "PNG")
        except Exception:
            output_path.write_bytes(image_bytes)
